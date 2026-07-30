#!/usr/bin/env python3
"""Generate Sustainability Dashboard data-flow documentation as PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_GREEN = RGBColor(0x1B, 0x5E, 0x20)
MID_GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
ACCENT = RGBColor(0xFF, 0x8F, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)

OUTPUT = "/workspace/docs/Sustainability-Dashboard-Data-Flow.pptx"


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_GREEN
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xC8, 0xE6, 0xC9)


def add_bullets(slide, items, left=0.6, top=1.4, width=12.0, height=5.8, font_size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
            p.text = text
            p.level = level
        else:
            p.text = item
            p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)


def add_table_slide(slide, headers, rows, top=1.5):
    cols = len(headers)
    tbl_shape = slide.shapes.add_table(
        len(rows) + 1, cols, Inches(0.5), Inches(top), Inches(12.3), Inches(0.35 * (len(rows) + 2))
    )
    tbl = tbl_shape.table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = MID_GREEN
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK_TEXT


def add_flow_box(slide, text, left, top, width=2.2, height=0.65, fill=MID_GREEN):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = DARK_GREEN
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_arrow_text(slide, text, left, top):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(1.2), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_GREEN)
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5))
    p = t.text_frame.paragraphs[0]
    p.text = "Sustainability Dashboard"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    s = t.text_frame.add_paragraph()
    s.text = "Dokumentasi Alur Data & Rumus Bisnis"
    s.font.size = Pt(24)
    s.font.color.rgb = RGBColor(0xC8, 0xE6, 0xC9)
    meta = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11), Inches(1.2))
    mp = meta.text_frame.paragraphs[0]
    mp.text = "Mill Onboarding · Supplier Due Diligence · Modul Terkait"
    mp.font.size = Pt(16)
    mp.font.color.rgb = RGBColor(0xA5, 0xD6, 0xA7)
    mp2 = meta.text_frame.add_paragraph()
    mp2.text = "Sumber: codebase valid (main.js, GoogleAppsScript-backend-v3-full.gs)"
    mp2.font.size = Pt(12)
    mp2.font.color.rgb = RGBColor(0x81, 0xC7, 0x84)


def content_slide(prs, title, subtitle, bullets, font_size=16):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, title, subtitle)
    add_bullets(slide, bullets, font_size=font_size)


def table_slide(prs, title, subtitle, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, title, subtitle)
    add_table_slide(slide, headers, rows)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)

    # 2 Architecture
    content_slide(
        prs,
        "Arsitektur Sistem",
        "Frontend → API Proxy → Google Apps Script → Google Sheets",
        [
            "Frontend: Vite SPA (src/main.js + modul panel)",
            "API: /api/gas-proxy → doGet / doPost di Google Apps Script",
            "Database: Google Sheets (bukan SQL) — setiap modul = 1+ tab sheet",
            "Prinsip kunci:",
            ("Skor/risk mill dihitung di Google Sheet (formula), bukan di web app", 1),
            ("Web app TIDAK menulis ulang kolom formula saat save", 1),
            ("Backend canonical: scripts/GoogleAppsScript-backend-v3-full.gs", 1),
        ],
    )

    # 3 Architecture flow visual
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, "Arsitektur — Diagram Alur", "Semua panel berbagi apiGet / apiPost")
    boxes = [
        ("Supplier\nDue Diligence", 0.5, 1.6, LIGHT_GREEN, DARK_TEXT),
        ("Mill\nOnboarding", 2.9, 1.6, LIGHT_GREEN, DARK_TEXT),
        ("Monitoring\nTTM/TTP", 5.3, 1.6, LIGHT_GREEN, DARK_TEXT),
        ("EUDR\nPotential", 7.7, 1.6, LIGHT_GREEN, DARK_TEXT),
        ("Monthly\nReport", 10.1, 1.6, LIGHT_GREEN, DARK_TEXT),
    ]
    for text, l, t, fill, txt in boxes:
        sh = add_flow_box(slide, text, l, t, 2.0, 0.9, MID_GREEN if txt == WHITE else fill)
        for p in sh.text_frame.paragraphs:
            p.font.color.rgb = txt if txt != WHITE else WHITE
            p.font.size = Pt(12)
    add_flow_box(slide, "/api/gas-proxy", 4.5, 3.0, 4.0, 0.7, ACCENT)
    add_flow_box(slide, "Google Apps Script\n(doGet / doPost)", 3.8, 4.0, 5.2, 0.9, DARK_GREEN)
    add_flow_box(slide, "Google Sheets\n(SDD_MAIN, Mill Profile, TTP, ...)", 3.2, 5.2, 6.4, 0.9, MID_GREEN)
    for l in [1.5, 3.9, 6.3, 8.7, 11.1]:
        add_arrow_text(slide, "↓", l, 2.55)

    # 4 SDD overview
    content_slide(
        prs,
        "Supplier Due Diligence (SDD)",
        "Panel: #panel-supplier-dd · Penyimpanan: SDD_MAIN, SDD_MILL_LIST, SDD_FFB_LIST",
        [
            "Entry point:",
            ("Upload Excel (Main Form + Traceability A/B)", 1),
            ("Screening form manual (SCR fields, GRV, PRI)", 1),
            ("Approve / Hold / Reject oleh approver", 1),
            ("NBL Check vs sheet NBL & Unilever NBL", 1),
            "Primary key: submission_id (format SUB-YYYYMMDD-XXXXXX)",
            "Supplier type: MILL · KCP · TRADER",
            "Legacy: sheet SDD Data (flat) masih didukung via adapter",
        ],
    )

    # 5 SDD state machine
    content_slide(
        prs,
        "SDD — State Machine",
        "Lifecycle screening dari import hingga onboarding selesai",
        [
            "Import Excel → Save Draft (SCR Status = Draft)",
            "Submit → SCR Status = Submitted (form terkunci, decision enabled)",
            "Boss Decision:",
            ("Approve → statusSDD = APPROVED", 2),
            ("On Hold → statusSDD = ON HOLD", 2),
            ("Reject → statusSDD = REJECTED", 2),
            "Jika APPROVED + Submitted → muncul di Mill Onboarding Task List",
            "User klik '+ Add to Mill' → data masuk Mill Onboarding Profile",
            "Setelah semua TML line selesai → mill_added = true (hilang dari Task List)",
            "TTP sync terjadi SETELAH Mill disimpan (bukan saat approve SDD saja)",
        ],
    )

    # 6 SDD downstream
    table_slide(
        prs,
        "SDD → Kemana Data Mengalir?",
        "Trigger, target, dan fungsi backend",
        ["Trigger", "Target Modul", "Fungsi / Mekanisme"],
        [
            ["statusSDD = APPROVED", "Contact List Supplier", "syncContactFromApprovedSubmission_()"],
            ["APPROVED + Submitted", "Mill Task List", "isApprovedSubmittedSddRow_() filter UI"],
            ["Save Mill dari Task List", "Mill Onboarding Profile", "add / update sheet mill"],
            ["mill_ttp_sync payload", "Monitoring TTP/TTM", "syncTtpFromMillOnboarding_*()"],
            ["NBL Check", "Kolom SCR - No Buy List", "runSddNblCheck() vs NBL sheets"],
            ["Monthly Report", "Section 01 SDD", "fetchSddList() / apiListSubmissions"],
            ["Export PDF", "SDD Screening PDF", "sddExportPdf()"],
        ],
    )

    # 7 SDD field mapping
    table_slide(
        prs,
        "Mapping Field: SDD → Mill Onboarding",
        "mapSddRowToMillPayload() / mapSddTmlRowToMillPayload()",
        ["Kolom Mill", "Sumber SDD"],
        [
            ["GROUP NAME", "Group Name"],
            ["COMPANY NAME", "Company Name"],
            ["MILL NAME", "Mill Name / KCP Name / TML - Mill Name"],
            ["COORDINATES", "Latitude + Longitude (comma-separated)"],
            ["MILL CAPACITY (TON/HOUR)", "Mill Capacity (Ton/Hour)"],
            ["MILL CATEGORY", "Mill Category / KCP Category"],
            ["ADDRESS", "Mill Address / KCP Address"],
            ["MONTH, YEAR", "Periode SDD (resolveSddPeriodFromMainRow_)"],
            ["SOURCE TYPE", "MILL / KCP / TRADER (millSourceTypeForOnboarding_)"],
            ["TRADER NAME", "Nama trader parent (jika applicable)"],
            ["UML ID", "TML - UML ID (untuk TML lines)"],
        ],
    )

    # 8 SDD FFB to TTP
    table_slide(
        prs,
        "Mapping: SDD FFB → Monitoring TTP/TTM",
        "1 baris TTP = 1 pemasok FFB (Traceability B) + konteks Main Form",
        ["Kolom TTP", "Sumber SDD"],
        [
            ["GROUP / COMPANY / MILL NAME", "Mill Onboarding identity (setelah Task List save)"],
            ["FFB SUPPLIER GROUP/NAME", "FFB - Supplier Group/Name"],
            ["CATEGORY", "FFB - Supplier Category"],
            ["LAT, LONG", "FFB - Latitude/Longitude"],
            ["VILLAGE, SUBDISTRICT, DISTRICT", "FFB location fields"],
            ["CONCESION / PLANTED AREA", "FFB - Concession/Planted Area (Ha)"],
            ["LEGALITAS, ISPO/RSPO/ISCC", "FFB certification fields"],
            ["FFB SUPPLY to MILL (TON)", "FFB - Total Supply FFB (Ton)"],
            ["submission_id, ffb_line_id", "Metadata sync (upsert identifier)"],
        ],
    )

    # 9 Mill onboarding entry
    content_slide(
        prs,
        "Mill Onboarding — Entry Point",
        "Sheet: Mill Onboarding Profile (CPO/PK) · Mill Onboarding Waste (POME/SHELL)",
        [
            "3 cara data masuk:",
            ("Manual Add/Edit — modal #btn-add-mill", 1),
            ("SDD Task List — '+ Add to Mill' dari approved screening", 1),
            ("Supply Excel Import — upload → draft → submit", 1),
            "Supply routing:",
            ("CPO/PK → sheet mill", 2),
            ("POME (ISCC/INS) / SHELL → sheet millWaste", 2),
            ("KCP plant → FACILITY NAME PK", 2),
            ("Refinery plant → FACILITY NAME CPO", 2),
            "Staging supply: Supply Import Draft (sheet supplyDraft)",
        ],
    )

    # 10 Mill outbound
    table_slide(
        prs,
        "Mill Onboarding → Kemana Data Mengalir?",
        "Consumer modules dan join key",
        ["Consumer", "Join Key", "Data yang Dipakai"],
        [
            ["Monitoring TTP/TTM", "GROUP|COMPANY|MILL|UML ID", "Identity + FFB sync dari SDD"],
            ["Monthly Report §02", "Period + RESULT RISK LEVEL", "High-risk MILL-only rows"],
            ["Performa Facility", "FACILITY NAME CPO/PK", "Supply qty, traceability %, ISPO"],
            ["Questionnaire Monitoring", "group|company|mill", "Registry mill list"],
            ["EUDR Potential", "group|company|mill", "Identity sync (STATUS manual)"],
            ["Mill Executive Report", "Quarter aggregation", "Supply, risk, NBL, cert stats"],
            ["Risk Reason UI", "Per-row gaps", "RESULT RISK LEVEL + field gaps"],
        ],
    )

    # 11 Supply task list
    content_slide(
        prs,
        "Supply Task List — Alur",
        "Embedded di Mill Onboarding panel",
        [
            "Excel upload → parse (PLANT, CATEGORY, COMPANY, MILL, QTY SUPPLY)",
            "Match ke profil mill: supplyFindMillProfileMatch_()",
            ("matched — profil ditemukan", 1),
            ("new — profil baru", 1),
            ("group_mismatch — perlu review", 1),
            "Save Draft → sheet Supply Import Draft",
            "Submit → submitSupplyDraft_ → append row ke Mill/Waste sheet",
            "Formula columns di-restore dari baris di atas (tidak di-overwrite)",
            "Repair: reconcileSupplyDraft_() perbaiki drift status draft vs mill",
        ],
    )

    # 12 Sheet formulas
    content_slide(
        prs,
        "Rumus di Google Sheet (Read-Only dari Web)",
        "Kolom MILL_FORMULA_HEADERS_ — tidak pernah ditulis dari web app",
        [
            "Skor: SCORE, DEFORESTATION/BURN/PEAT SCORE, TOTAL SCORE SPATIAL, TOTAL SCORE, LEGALITY SCORE",
            "Agregat: TOTAL GRIEVANCES, TOTAL POLICY, TOTAL CERTIFICATION",
            "Risk: RISK LEVEL, RESULT RISK LEVEL, BUYER NO BUY LIST, PRIORITY ENGAGEMENT, SUPPLIER STATUS",
            "Supply %: PERCENTAGE SUPPLY CPO/PK/ISCC/INS/SHELL, PRODUCT SUPPLY, STATUS SUPPLY CPO/PK",
            "Waste: TOTAL POME SUPPLY, MAX SUPPLY POME/SHELL, REMAINING STOCK, TOTAL SCORE SUPPLY, TOTAL RISK LEVEL",
            "Status: COMPLIMENT/NOT COMPLIMENT, DECLARATION MONITORING, CPO, PK",
            "SD Monitoring: LAST SD DATE, DAY LEFT, Status, Result, Risk Number",
        ],
        font_size=14,
    )

    # 13 Grievance formula
    content_slide(
        prs,
        "Rumus App: Grievance Risk Score",
        "src/grievance-risk.js — dihitung di frontend, disimpan ke sheet",
        [
            "Total Score = Σ 6 indikator (masing-masing skor 1–3):",
            ("Publish Grievance: Yes=2, No=1", 1),
            ("Subject: Group Level=2, Subsidiary=1", 1),
            ("Repeat: 2×+=3, 1×=2, Not repeated=1", 1),
            ("Consequence: Harm=3, Action plan=2, No impact=1", 1),
            ("Group Scale: 5+ mills=3, 2–4=2, 1 mill=1", 1),
            ("No Response: No response 2 weeks=2, Response=1", 1),
            "Klasifikasi:",
            ("Total ≤ 8 → Low", 1),
            ("Total ≤ 11 → Medium", 1),
            ("Total > 11 → High", 1),
        ],
    )

    # 14 TTM TTP formulas
    content_slide(
        prs,
        "Rumus App: TTM & TTP Traceability %",
        "src/main.js — ttpCalcTtmCoordinatePct_ / ttpAggregateTotalTraceablePct_",
        [
            "TTM Coordinate %:",
            ("Pool = SOURCE TYPE ∈ {MILL, TRADER, REFINERY} AND SUPPLY CPO/PK > 0", 1),
            ("Traceable = baris dengan koordinat valid (lat/lng)", 1),
            ("% = Σ(qty traceable) / Σ(qty total) × 100", 1),
            "TTP Aggregate %:",
            ("Primary: Σ(Traceable Volume) / Σ(CPO/PK SUPPLY to REFINERY/KCP) × 100", 1),
            ("Fallback: average kolom % CPO/PK TRACEABLE", 1),
            "Performa Facility (weighted):",
            ("TTM Facility = Σ(qty × ttmPct/100) / Σ(qty) × 100", 1),
            ("ISPO Supply % = SUM(ISPO qty) / totalQty × 100", 1),
        ],
    )

    # 15 EUDR formula
    content_slide(
        prs,
        "Rumus App: EUDR Potential Status",
        "eudrComputeStatus_() — semua kriteria ENABLED harus PASS",
        [
            "Semua pass → Potential · Salah satu fail → Not Potential",
            "Kriteria (configurable via EUDR Status Formula sheet):",
            ("legality — Legality Complete / score = 1", 1),
            ("millCategory — contains 'Integrated'", 1),
            ("ownPlasmaFfb — FFB % per kategori TTP ≥ threshold (default 70%)", 1),
            ("resultRiskLevel — RESULT RISK LEVEL = Low", 1),
            ("millLocation — MILL LOC contains 'APL'", 1),
            ("certification — ≥1 RSPO/ISPO/ISCC = Yes", 1),
            ("grievance — count ≤ 0 · ndpePolicy — NDPE = Yes", 1),
            ("noBuyList — not on NBL · deforestation — Ha < 10 ha OR After 2020 = No", 1),
        ],
        font_size=14,
    )

    # 16 Mill risk reason
    content_slide(
        prs,
        "Rumus App: Mill Risk Reason Pills",
        "src/mill-risk-reason.js — interpretasi gap field untuk UI display",
        [
            "No Coordinate — koordinat kosong/invalid",
            "Legality Not Complete — LEGALITY SCORE ≠ 1",
            "Non APL Area — MILL LOC bukan APL",
            "No NDPE — NDPE ≠ Yes · No Certification — tidak ada cert",
            "Deforestation — RISK REDUCTION FACTOR = 1 OR width > 0 ha",
            "High Deforestation — factor = 2 OR width > 25 ha",
            "On No Buy List — BUYER NO BUY LIST = Yes",
            "Grievance flags — LEGALITY/ENVIRONMENT GRIEVANCE = Yes",
            "Display: uses RESULT RISK LEVEL; LOW mills tidak tampilkan reason pills",
        ],
        font_size=14,
    )

    # 17 Module map visual
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, "Peta Modul & Koneksi", "Alur dependensi antar modul")
    flows = [
        ("SDD", 0.4, 1.5), ("Mill", 3.0, 1.5), ("TTP", 5.6, 1.5),
        ("EUDR", 8.2, 1.5), ("Monthly Report", 10.4, 1.5),
        ("Contact", 0.4, 3.2), ("Performa Facility", 3.0, 3.2),
        ("Questionnaire", 5.6, 3.2), ("Grievance", 8.2, 3.2), ("NBL", 10.4, 3.2),
    ]
    for text, l, t in flows:
        add_flow_box(slide, text, l, t, 2.2, 0.65, MID_GREEN)
    notes = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.8))
    tf = notes.text_frame
    tf.word_wrap = True
    lines = [
        "SDD (APPROVED) → Contact List · SDD → Mill Task List → Mill Profile → TTP sync",
        "Supply Excel → Mill Profile · Mill → EUDR Potential · TTP FFB → EUDR FFB %",
        "Mill + TTP + Grievance + NBL + EUDR → Monthly Report",
        "Mill + TTP + Supplied CPO/PK → Performa Facility",
        "NBL → SDD screening check + Mill BUYER NO BUY LIST (sheet formula)",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "→ " + line
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT

    # 18 Sheet registry
    table_slide(
        prs,
        "Daftar Google Sheets (Database)",
        "Semua tab sheet yang digunakan sistem",
        ["Key", "Tab Name", "Modul"],
        [
            ["mill", "Mill Onboarding Profile", "Mill Onboarding"],
            ["millWaste", "Mill Onboarding Waste", "Mill Onboarding Waste"],
            ["supplyDraft", "Supply Import Draft", "Supply Task List"],
            ["sddMain / sddMill / sddFfb", "SDD_MAIN / MILL_LIST / FFB_LIST", "Supplier Due Diligence"],
            ["ttp", "Monitoring TTP/TTM", "TTM/TTP"],
            ["grievance", "Grievance Monitoring", "Grievance"],
            ["nbl / unileverNbl", "NBL / Unilever NBL", "No Buy List"],
            ["eudrPotential", "EUDR Potential", "EUDR Potential"],
            ["eudrDds (+ child)", "EUDR DDS + Suppliers/Geo/Docs", "EUDR DDS (export EU)"],
            ["contactSupplier", "Contact List Supplier", "Contact List"],
            ["questionnaireMonitoring", "Questionnaire Monitoring", "Questionnaire"],
        ],
    )

    # 19 SDD vs DDS
    table_slide(
        prs,
        "Perbedaan: SDD vs EUDR DDS",
        "Dua modul 'due diligence' yang terpisah — tidak ada coupling langsung",
        ["Aspek", "Supplier Due Diligence (SDD)", "EUDR Due Diligence Statement (DDS)"],
        [
            ["Panel", "#panel-supplier-dd", "#panel-due-diligence-statement"],
            ["Storage", "SDD_MAIN / MILL_LIST / FFB_LIST", "EUDR DDS + child sheets"],
            ["Tujuan", "Screening supplier internal + onboarding", "Paket data export EU EUDR"],
            ["Output", "Task List → Mill → TTP", "DOCX/PDF export untuk buyer"],
            ["Code", "src/main.js", "src/dds-ui.js, dds-export-model.js"],
        ],
    )

    # 20 Validation
    content_slide(
        prs,
        "Aturan Validasi & Guard Rules",
        "Business rules yang mencegah data inconsistency",
        [
            "SDD: Submitted tidak bisa downgrade ke Draft (assertSubmittedNotDowngraded_)",
            "Task List gate: SCR Status = submitted AND statusSDD = APPROVED",
            "TTP sync: APPROVED + submitted + identity lengkap (GROUP/COMPANY/MILL/UML)",
            "Supply submit: match_status harus matched atau new; company wajib diisi",
            "Mill save: kolom formula di-strip (millStripComputedFromSavePayload_)",
            "EUDR: semua kriteria enabled harus pass untuk status Potential",
            "Grievance: 6 indikator wajib diisi sebelum total/classification disimpan",
            "Koordinat: normalizeCoordinate() + recoverCoord() untuk locale Indonesia",
        ],
        font_size=14,
    )

    # 21 Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_GREEN)
    t = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    p = t.text_frame.paragraphs[0]
    p.text = "Terima Kasih"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    s = t.text_frame.add_paragraph()
    s.text = "Dokumentasi valid dari codebase · Tidak ada perubahan kode aplikasi"
    s.font.size = Pt(16)
    s.font.color.rgb = RGBColor(0xC8, 0xE6, 0xC9)
    s.alignment = PP_ALIGN.CENTER
    ref = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1.5))
    rp = ref.text_frame.paragraphs[0]
    rp.text = "Referensi: docs/sdd-to-ttm-ttp-field-mapping.md"
    rp.font.size = Pt(12)
    rp.font.color.rgb = RGBColor(0xA5, 0xD6, 0xA7)
    rp.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build_presentation()
